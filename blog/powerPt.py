from pptx.enum.dml import MSO_FILL_TYPE
from pptx.util import Inches, Pt
from .models import BlogPost
from django.shortcuts import render, get_object_or_404
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from PIL import Image


def blog_post_powerpoint(request, slug):
    obj = get_object_or_404(BlogPost, slug=slug)
    template_name = 'blog/retrieve.html'
    prs = Presentation('powerpoint/model.pptx')

    custom_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(custom_slide_layout)

    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))

    title = slide.shapes.title
    title_tf = title.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    p = title_tf.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    run = p.add_run()
    run.text = obj.title

    font = run.font
    font.size = Pt(50)
    font.bold = True

    p_content = title_tf.add_paragraph()
    p = title_tf.paragraphs[1]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    run = p.add_run()
    run.text = obj.content

    font = run.font
    font.size = Pt(20)
    font.bold = True

    img_path = 'powerpoint/altran.jpg'
    picture_placeholder = slide.placeholders[11]
    placeholder_picture = picture_placeholder.insert_picture(img_path)
    placeholder_picture.crop_top = 0
    placeholder_picture.crop_bottom = 0

    subtitle_placeholder = slide.placeholders[13]
    subtitle_tf = subtitle_placeholder.text_frame
    subtitle_tf.word_wrap = True
    subtitle_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    subtitle_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    p = subtitle_tf.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    run = p.add_run()
    run.text = "Defect Image"

    font = run.font
    font.size = Pt(20)
    font.bold = False

    for item in prs.slides:
        print(item.slide_id)

    if obj.image:
        img_path = obj.image
        picture_placeholder = slide.placeholders[16]
        placeholder_picture = picture_placeholder.insert_picture(img_path)
        placeholder_picture.crop_top = 0
        placeholder_picture.crop_bottom = 0

        subtitle_placeholder = slide.placeholders[17]
        subtitle_tf = subtitle_placeholder.text_frame
        subtitle_tf.word_wrap = True
        subtitle_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        subtitle_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        p = subtitle_tf.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        run = p.add_run()
        run.text = "Defect Image"

        font = run.font
        font.size = Pt(20)
        font.bold = False

    table_placeholder_first = slide.placeholders[10]
    graphic_frame_first = table_placeholder_first.insert_table(rows=5, cols=9)
    table_first = graphic_frame_first.table

    for row in table_first.rows:
        row.height = Inches(0.5)

    first_row_first = table_first.rows[0]
    first_row_first.height = Inches(1)

    table_placeholder_second = slide.placeholders[15]
    graphic_frame_second = table_placeholder_second.insert_table(rows=3, cols=5)
    table_second = graphic_frame_second.table
    table_second.fixed_width = True

    for row in table_second.rows:
        row.height = Inches(1.5)

    first_row_second = table_second.rows[0]
    first_row_second.height = Inches(0.4)

    title_cell = table_second.cell(0, 0)
    title_cell.merge(table_second.cell(0, 3))

    table_second.columns[0].width = Inches(1)
    table_second.columns[1].width = Inches(4)
    table_second.columns[2].width = Inches(1)
    table_second.columns[3].width = Inches(3.7)
    table_second.columns[4].width = Inches(3.6)

    img_path = 'powerpoint/Camembert.png'
    picture_cell = table_second.cell(0, 1)

   

    # write column headings
    # table.cell(0, 0).text = 'Foo'
    # table.cell(0, 1).text = 'Bar'

    # write body cells
    # table.cell(1, 0).text = 'Baz'
    # table.cell(1, 1).text = 'Qux'

    prs.save('powerpoint/test.pptx')
    context = {"object": obj, "presentation": prs}
    return render(request, template_name, context)
